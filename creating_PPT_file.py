from pptx import Presentation
from pptx.util import Inches, Pt

prs = Presentation()

def add_slide(title_text, body_text):
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    content = slide.placeholders[1]
    title.text = title_text
    content.text = body_text

add_slide("1. Identify the Social Problem", 
          "• Issue: Systemic ethnic conflict and humanitarian crisis.\n"
          "• Affected: Kuki-Zo and Meitei communities in Manipur.\n"
          "• Manifestation: Armed skirmishes, arson, and total physical segregation.\n"
          "• Observation: Collapse of law and order following leadership vacuums.")
add_slide("2. Collect Relevant Data", 
          "• Deaths: 260+ confirmed fatalities.\n"
          "• Displacement: 60,000+ people in relief camps.\n"
          "• Power Dynamics: Imposition of President's Rule (Feb 2025).\n"
          "• Gaps: Under-reporting of sexual violence and missing persons.")

add_slide("3. Social Pathology & Root Causes", 
          "• Social: Ethnic 'buffer zones' replacing community networks.\n"
          "• Economic: 300% inflation due to highway blockades.\n"
          "• Political: Partisanship in law enforcement; 'Scheduled Tribe' status trigger.\n"
          "• Cultural: Dehumanizing narratives (e.g., 'illegal immigrants').")

add_slide("4. Identify Stakeholders", 
          "• Impacted: Displaced families, students, and women survivors.\n"
          "• Contributors: Armed vigilante groups and polarizing political rhetoric.\n"
          "• Solvers: Central Government (MHA), Supreme Court, and NGOs (HRW).")

add_slide("5. Impact & Ethical Reflection", 
          "• Long-term: Generational trauma and permanent ethnic enclaves.\n"
          "• Inequality: Disproportionate suffering of minorities in both regions.\n"
          "• Ethics: Responsibility of the state to maintain impartial protection.")


add_slide("6. Develop Possible Solutions", 
          "• Preventive: Neutral Security Corridors.\n"
          "• Corrective: Fast-track courts for ethnic violence cases.\n"
          "• Policy: Truth and Reconciliation Commission.\n"
          "• Tech: AI-driven hate speech monitoring.")
add_slide("7. Solemn Pledge", 
          "• Action: I pledge to verify all news regarding communal tension.\n"
          "• Why: To stop the spread of misinformation that fuels violence.\n"
          "• Practice: Using fact-checking tools and challenging dehumanizing language.")

prs.save('Manipur_Social_Analysis.pptx')
print("Presentation created successfully!")